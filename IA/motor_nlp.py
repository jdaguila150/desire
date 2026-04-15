from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
import os

# Asegúrate de importar tu motor de imágenes
from IA.motor_imagenes import descargar_imagen_unsplash 

TEMAS = {
    "corporativo": {"acento": RGBColor(0, 82, 165), "texto_titulo": RGBColor(0, 51, 102)},
    "creativo": {"acento": RGBColor(255, 102, 0), "texto_titulo": RGBColor(200, 80, 0)},
    "oscuro": {"acento": RGBColor(0, 200, 150), "texto_titulo": RGBColor(30, 30, 30)}
}

def aplicar_capa_diseno(slide, prs, color_acento):
    barra = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), prs.slide_height)
    barra.fill.solid()
    barra.fill.fore_color.rgb = color_acento
    barra.line.fill.background()

def crear_presentacion_desde_json(datos_json, nombre_archivo="Presentacion_Final.pptx"):
    prs = Presentation()
    tema_elegido = datos_json.get("tema", "corporativo")
    paleta = TEMAS.get(tema_elegido, TEMAS["corporativo"])

    # --- 1. PORTADA ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    aplicar_capa_diseno(slide, prs, paleta["acento"])
    slide.shapes.title.text = datos_json.get("titulo_principal", "Desire Presentation")
    slide.placeholders[1].text = datos_json.get("subtitulo", "Generado con IA")

    # --- 2. CICLO DE DIAPOSITIVAS ---
    for diapo in datos_json.get("diapositivas", []):
        layout = diapo.get("layout", "clasico")
        titulo = diapo.get("titulo", "Sin Título")
        lista_puntos = diapo.get("contenido", diapo.get("puntos", []))
        texto_insight = diapo.get("insight", diapo.get("dato_extra", ""))

        query_imagen = diapo.get("query_imagen", "")
        ruta_img_local = descargar_imagen_unsplash(query_imagen)

        # Usaremos layout 5 (Solo título) para tener un lienzo limpio y control total
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        aplicar_capa_diseno(slide, prs, paleta["acento"])

        # Control del Título
        title_shape = slide.shapes.title
        title_shape.text = titulo
        title_shape.text_frame.paragraphs[0].font.color.rgb = paleta["texto_titulo"]

        # --- DISEÑO DESTACADO ---
        if layout == "destacado":
            ancho_texto = Inches(4) if ruta_img_local else Inches(8)
            text_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), ancho_texto, Inches(3))
            tf = text_box.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

            p = tf.paragraphs[0]
            p.text = texto_insight if texto_insight else titulo
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = paleta["texto_titulo"]

            if ruta_img_local:
                slide.shapes.add_picture(ruta_img_local, Inches(5.2), Inches(1.5), width=Inches(4.3))

        # --- DISEÑO INSIGHT ---
        elif layout == "insight":
            ancho_texto = Inches(4.2) if ruta_img_local else Inches(8)
            
            # Dibujamos nuestra propia caja de texto EXACTAMENTE donde la queremos
            caja_puntos = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), ancho_texto, Inches(3.2))
            tf = caja_puntos.text_frame
            tf.word_wrap = True

            for i, pto in enumerate(lista_puntos):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = pto
                p.font.size = Pt(22)
                p.level = 0               # <--- FUERZA LA VIÑETA
                p.space_after = Pt(14)    # <--- EVITA QUE SE ENCIMEN

            if ruta_img_local:
                slide.shapes.add_picture(ruta_img_local, Inches(5.2), Inches(1.5), width=Inches(4.3))

            if texto_insight:
                caja = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.3), Inches(9), Inches(1.5))
                caja.fill.solid()
                caja.fill.fore_color.rgb = paleta["acento"]
                caja.line.fill.background()

                tf_ins = caja.text_frame
                p_ins = tf_ins.paragraphs[0]
                p_ins.text = f"💡 {texto_insight}"
                p_ins.font.size = Pt(22)
                p_ins.font.bold = True
                p_ins.font.color.rgb = RGBColor(255, 255, 255)

        # --- DISEÑO CLÁSICO ---
        else:
            ancho_texto = Inches(4.2) if ruta_img_local else Inches(8)
            
            # Dibujamos nuestra propia caja de texto
            caja_puntos = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), ancho_texto, Inches(4))
            tf = caja_puntos.text_frame
            tf.word_wrap = True

            for i, pto in enumerate(lista_puntos):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = pto
                p.font.size = Pt(24)
                p.level = 0               # <--- FUERZA LA VIÑETA
                p.space_after = Pt(14)    # <--- EVITA QUE SE ENCIMEN

            if ruta_img_local:
                slide.shapes.add_picture(ruta_img_local, Inches(5.2), Inches(1.8), width=Inches(4.3))

    prs.save(nombre_archivo)
    return nombre_archivo